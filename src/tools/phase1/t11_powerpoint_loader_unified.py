"""
T11: PowerPoint Document Loader - Unified Interface Implementation

Loads and parses PowerPoint documents (.pptx, .ppt) with comprehensive content extraction.
"""

from typing import Dict, Any, Optional, List, Union
import os
from pathlib import Path
import uuid
from datetime import datetime
import logging

from src.tools.base_tool import BaseTool, ToolRequest, ToolResult, ToolContract, ToolStatus, ToolErrorCode
from src.core.service_manager import ServiceManager
from src.core.advanced_data_models import Document, ObjectType, QualityTier

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


class T11PowerPointLoaderUnified(BaseTool):
    """T11: PowerPoint Document Loader with unified interface"""
    
    def __init__(self, service_manager: ServiceManager):
        """Initialize with service manager"""
        super().__init__(service_manager)
        self.tool_id = "T11_POWERPOINT_LOADER"
        self.identity_service = service_manager.identity_service
        self.provenance_service = service_manager.provenance_service
        self.quality_service = service_manager.quality_service
        self._temp_files = []
    
    def get_contract(self) -> ToolContract:
        """Return tool contract specification"""
        return ToolContract(
            tool_id=self.tool_id,
            name="PowerPoint Document Loader",
            description="Load and parse PowerPoint documents with comprehensive content extraction",
            category="document_processing",
            input_schema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to PowerPoint file to load"
                    },
                    "workflow_id": {
                        "type": "string",
                        "description": "Optional workflow ID for tracking"
                    },
                    "parse_options": {
                        "type": "object",
                        "properties": {
                            "extract_images": {"type": "boolean", "default": False},
                            "extract_notes": {"type": "boolean", "default": True},
                            "extract_metadata": {"type": "boolean", "default": True},
                            "include_hidden_slides": {"type": "boolean", "default": False}
                        },
                        "default": {}
                    }
                },
                "required": ["file_path"]
            },
            output_schema={
                "type": "object",
                "properties": {
                    "document": {
                        "type": "object",
                        "properties": {
                            "document_id": {"type": "string"},
                            "document_ref": {"type": "string"},
                            "file_path": {"type": "string"},
                            "file_name": {"type": "string"},
                            "file_size": {"type": "integer"},
                            "presentation_data": {"type": "object"},
                            "text_content": {"type": "string"},
                            "slide_count": {"type": "integer"},
                            "shape_count": {"type": "integer"},
                            "confidence": {"type": "number"},
                            "quality_tier": {"type": "string"},
                            "created_at": {"type": "string"}
                        },
                        "required": ["document_id", "presentation_data", "text_content", "confidence", "slide_count"]
                    }
                },
                "required": ["document"]
            },
            dependencies=["identity_service", "provenance_service", "quality_service"],
            performance_requirements={
                "max_execution_time": 90.0,   # 90 seconds for large presentations
                "max_memory_mb": 1024,        # 1GB for PowerPoint processing
                "min_confidence": 0.8         # Minimum confidence threshold
            },
            error_conditions=[
                "FILE_NOT_FOUND",
                "INVALID_FILE_TYPE",
                "POWERPOINT_CORRUPTED",
                "POWERPOINT_PASSWORD_PROTECTED",
                "PPTX_LIBRARY_MISSING",
                "MEMORY_LIMIT_EXCEEDED"
            ]
        )
    
    def execute(self, request: ToolRequest) -> ToolResult:
        """Execute PowerPoint loading with unified interface"""
        self._start_execution()
        
        try:
            # Check if python-pptx is available
            if not PPTX_AVAILABLE:
                return self._create_error_result(
                    request,
                    "PPTX_LIBRARY_MISSING",
                    "python-pptx library is not installed. Install with: pip install python-pptx"
                )
            
            # Validate input
            if not self.validate_input(request.input_data):
                return self._create_error_result(
                    request,
                    "INVALID_INPUT",
                    "Input validation failed. Required: file_path"
                )
            
            # Extract parameters
            file_path = request.input_data.get("file_path")
            workflow_id = request.input_data.get("workflow_id")
            parse_options = request.input_data.get("parse_options", {})
            
            # Set default parse options
            parse_options = {
                "extract_images": parse_options.get("extract_images", False),
                "extract_notes": parse_options.get("extract_notes", True),
                "extract_metadata": parse_options.get("extract_metadata", True),
                "include_hidden_slides": parse_options.get("include_hidden_slides", False)
            }
            
            # Validate file path
            validation_result = self._validate_file_path(file_path)
            if not validation_result["valid"]:
                return self._create_error_result(
                    request,
                    validation_result["error_code"],
                    validation_result["error_message"]
                )
            
            file_path = Path(file_path)
            
            # Start provenance tracking
            operation_id = self.provenance_service.start_operation(
                tool_id=self.tool_id,
                operation_type="load_powerpoint_document",
                used={},
                parameters={
                    "file_path": str(file_path),
                    "workflow_id": workflow_id,
                    "parse_options": parse_options
                }
            )
            
            # Generate workflow ID if not provided
            if not workflow_id:
                workflow_id = f"wf_{uuid.uuid4().hex[:8]}"
            
            # Create document ID
            document_id = f"{workflow_id}_{file_path.stem}"
            document_ref = f"storage://document/{document_id}"
            
            # Parse PowerPoint document
            parsing_result = self._parse_powerpoint_document(file_path, parse_options)
            
            if parsing_result["status"] != "success":
                return self._create_error_result(
                    request,
                    parsing_result.get("error_code", "POWERPOINT_PARSE_ERROR"),
                    parsing_result["error"]
                )
            
            # Calculate confidence
            confidence = self._calculate_confidence(
                presentation_data=parsing_result["presentation_data"],
                slide_count=parsing_result["slide_count"],
                shape_count=parsing_result["shape_count"],
                file_size=file_path.stat().st_size
            )
            
            # Create document data
            document_data = {
                "document_id": document_id,
                "document_ref": document_ref,
                "file_path": str(file_path),
                "file_name": file_path.name,
                "file_size": file_path.stat().st_size,
                "presentation_data": parsing_result["presentation_data"],
                "text_content": parsing_result["text_content"],
                "slide_count": parsing_result["slide_count"],
                "shape_count": parsing_result["shape_count"],
                "confidence": confidence,
                "created_at": datetime.now().isoformat(),
                "tool_version": "1.0.0",
                "parse_options": parse_options
            }
            
            # Assess quality
            quality_result = self.quality_service.assess_confidence(
                object_ref=document_ref,
                base_confidence=confidence,
                factors={
                    "slide_count": min(1.0, parsing_result["slide_count"] / 50),
                    "shape_count": min(1.0, parsing_result["shape_count"] / 200),
                    "file_size": min(1.0, file_path.stat().st_size / (50 * 1024 * 1024)),
                    "content_richness": min(1.0, len(parsing_result["text_content"]) / 10000)
                },
                metadata={
                    "powerpoint_type": "presentation",
                    "parse_method": "python-pptx"
                }
            )
            
            if quality_result["status"] == "success":
                document_data["confidence"] = quality_result["confidence"]
                document_data["quality_tier"] = quality_result["quality_tier"]
            
            # Complete provenance
            self.provenance_service.complete_operation(
                operation_id=operation_id,
                outputs=[document_ref],
                success=True,
                metadata={
                    "slide_count": parsing_result["slide_count"],
                    "shape_count": parsing_result["shape_count"],
                    "text_length": len(parsing_result["text_content"]),
                    "confidence": document_data["confidence"]
                }
            )
            
            # Get execution metrics
            execution_time, memory_used = self._end_execution()
            
            # Create success result
            return ToolResult(
                tool_id=self.tool_id,
                status="success",
                data={
                    "document": document_data
                },
                metadata={
                    "operation_id": operation_id,
                    "workflow_id": workflow_id,
                    "parse_method": "python-pptx",
                    "parse_options": parse_options
                },
                execution_time=execution_time,
                memory_used=memory_used
            )
            
        except Exception as e:
            logger.error(f"Unexpected error in {self.tool_id}: {e}", exc_info=True)
            return self._create_error_result(
                request,
                "UNEXPECTED_ERROR",
                f"Unexpected error during PowerPoint loading: {str(e)}"
            )
    
    def _validate_file_path(self, file_path: str) -> Dict[str, Any]:
        """Validate file path for security and existence"""
        if not file_path:
            return {
                "valid": False,
                "error_code": "INVALID_INPUT",
                "error_message": "File path cannot be empty"
            }
        
        try:
            path = Path(file_path)
            
            # Check if path exists
            if not path.exists():
                return {
                    "valid": False,
                    "error_code": "FILE_NOT_FOUND",
                    "error_message": f"File not found: {file_path}"
                }
            
            # Check if it's a file
            if not path.is_file():
                return {
                    "valid": False,
                    "error_code": "INVALID_INPUT",
                    "error_message": f"Path is not a file: {file_path}"
                }
            
            # Check extension
            allowed_extensions = ['.pptx', '.ppt']
            if path.suffix.lower() not in allowed_extensions:
                return {
                    "valid": False,
                    "error_code": "INVALID_FILE_TYPE",
                    "error_message": f"Invalid file extension. Allowed: {allowed_extensions}"
                }
            
            # Basic security check - prevent path traversal
            if ".." in str(path) or str(path).startswith("/etc"):
                return {
                    "valid": False,
                    "error_code": "VALIDATION_FAILED",
                    "error_message": "Invalid file path"
                }
            
            return {"valid": True}
            
        except Exception as e:
            return {
                "valid": False,
                "error_code": "VALIDATION_FAILED",
                "error_message": f"Path validation failed: {str(e)}"
            }
    
    def _parse_powerpoint_document(self, file_path: Path, parse_options: Dict[str, Any]) -> Dict[str, Any]:
        """Parse PowerPoint document using python-pptx"""
        try:
            # Load presentation
            presentation = Presentation(file_path)
            
            # Extract presentation metadata
            metadata = {}
            if parse_options.get("extract_metadata", True):
                core_props = presentation.core_properties
                metadata = {
                    "title": getattr(core_props, 'title', None),
                    "author": getattr(core_props, 'author', None),
                    "subject": getattr(core_props, 'subject', None),
                    "keywords": getattr(core_props, 'keywords', None),
                    "category": getattr(core_props, 'category', None),
                    "comments": getattr(core_props, 'comments', None),
                    "created": getattr(core_props, 'created', None),
                    "modified": getattr(core_props, 'modified', None),
                    "last_modified_by": getattr(core_props, 'last_modified_by', None)
                }
                # Convert datetime objects to strings
                for key, value in metadata.items():
                    if hasattr(value, 'isoformat'):
                        metadata[key] = value.isoformat()
            
            # Extract slides data
            slides_data = []
            text_content_parts = []
            total_shape_count = 0
            
            for slide_idx, slide in enumerate(presentation.slides):
                slide_data = {
                    "slide_number": slide_idx + 1,
                    "shapes": [],
                    "notes": None,
                    "text_content": ""
                }
                
                # Extract shapes and text
                slide_text_parts = []
                shape_count = 0
                
                for shape in slide.shapes:
                    shape_count += 1
                    shape_info = {
                        "shape_type": str(shape.shape_type),
                        "text": None,
                        "has_text": False
                    }
                    
                    # Extract text from shape
                    if hasattr(shape, "text") and shape.text:
                        shape_info["text"] = shape.text
                        shape_info["has_text"] = True
                        slide_text_parts.append(shape.text)
                    
                    # Handle tables
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_text = self._extract_table_text(shape.table)
                        if table_text:
                            shape_info["table_text"] = table_text
                            slide_text_parts.append(table_text)
                    
                    slide_data["shapes"].append(shape_info)
                
                slide_data["shape_count"] = shape_count
                total_shape_count += shape_count
                
                # Extract slide notes
                if parse_options.get("extract_notes", True):
                    try:
                        if slide.notes_slide and slide.notes_slide.notes_text_frame:
                            notes_text = slide.notes_slide.notes_text_frame.text
                            if notes_text.strip():
                                slide_data["notes"] = notes_text
                                slide_text_parts.append(f"Notes: {notes_text}")
                    except:
                        pass  # Notes extraction might fail for some slides
                
                # Combine slide text
                slide_text = " ".join(slide_text_parts)
                slide_data["text_content"] = slide_text
                
                if slide_text.strip():
                    text_content_parts.append(f"Slide {slide_idx + 1}: {slide_text}")
                
                slides_data.append(slide_data)
            
            # Combine all text content
            text_content = "\n\n".join(text_content_parts)
            
            # Create presentation data
            presentation_data = {
                "metadata": metadata,
                "slides": slides_data,
                "slide_count": len(presentation.slides),
                "total_shape_count": total_shape_count
            }
            
            return {
                "status": "success",
                "presentation_data": presentation_data,
                "text_content": text_content,
                "slide_count": len(presentation.slides),
                "shape_count": total_shape_count
            }
            
        except Exception as e:
            error_msg = str(e)
            if "password" in error_msg.lower() or "protected" in error_msg.lower():
                return {
                    "status": "error",
                    "error": f"PowerPoint file is password protected: {error_msg}",
                    "error_code": "POWERPOINT_PASSWORD_PROTECTED"
                }
            else:
                logger.error(f"Failed to parse PowerPoint document: {error_msg}")
                return {
                    "status": "error",
                    "error": f"Failed to parse PowerPoint document: {error_msg}",
                    "error_code": "POWERPOINT_CORRUPTED"
                }
    
    def _extract_table_text(self, table) -> str:
        """Extract text from a PowerPoint table"""
        try:
            table_text_parts = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text:
                        row_text.append(cell.text.strip())
                if row_text:
                    table_text_parts.append(" | ".join(row_text))
            return "\n".join(table_text_parts)
        except:
            return ""
    
    def _calculate_confidence(self, presentation_data: Dict[str, Any], slide_count: int, 
                            shape_count: int, file_size: int) -> float:
        """Calculate confidence score for PowerPoint parsing"""
        base_confidence = 0.9  # High confidence for successful PowerPoint parsing
        
        # Factors that affect confidence
        factors = []
        
        # Slide count factor
        if slide_count > 20:
            factors.append(0.95)
        elif slide_count > 5:
            factors.append(0.9)
        elif slide_count > 1:
            factors.append(0.85)
        else:
            factors.append(0.8)
        
        # Shape count factor (content richness)
        if shape_count > 100:
            factors.append(0.95)
        elif shape_count > 20:
            factors.append(0.9)
        elif shape_count > 5:
            factors.append(0.85)
        else:
            factors.append(0.8)
        
        # File size factor
        if file_size > 10 * 1024 * 1024:  # > 10MB
            factors.append(0.95)
        elif file_size > 1024 * 1024:  # > 1MB
            factors.append(0.9)
        else:
            factors.append(0.85)
        
        # Content quality factor
        has_meaningful_content = False
        for slide in presentation_data.get("slides", []):
            if slide.get("text_content", "").strip():
                has_meaningful_content = True
                break
        
        if has_meaningful_content:
            factors.append(0.95)
        else:
            factors.append(0.7)
        
        # Calculate average
        if factors:
            final_confidence = sum([base_confidence] + factors) / (len(factors) + 1)
        else:
            final_confidence = base_confidence
        
        # Ensure confidence is in valid range
        return max(0.1, min(1.0, final_confidence))
    
    def validate_input(self, input_data: Any) -> bool:
        """Validate input against tool contract with PowerPoint-specific validation"""
        # Call base validation first
        if not super().validate_input(input_data):
            return False
        
        # Additional validation for PowerPoint loader
        if isinstance(input_data, dict):
            file_path = input_data.get("file_path")
            if not file_path or not file_path.strip():
                return False
        
        return True
    
    def health_check(self) -> ToolResult:
        """Check tool health and readiness"""
        try:
            # Check if python-pptx is available
            if PPTX_AVAILABLE:
                from pptx import __version__ as pptx_version
                pptx_available = True
            else:
                pptx_version = 'not_installed'
                pptx_available = False
        except:
            pptx_available = False
            pptx_version = 'unknown'
        
        # Check service dependencies
        services_healthy = True
        if self.services:
            try:
                # Basic check that services exist
                _ = self.identity_service
                _ = self.provenance_service
                _ = self.quality_service
            except:
                services_healthy = False
        
        healthy = pptx_available and services_healthy
        
        return ToolResult(
            tool_id=self.tool_id,
            status="success" if healthy else "error",
            data={
                "healthy": healthy,
                "pptx_available": pptx_available,
                "pptx_version": pptx_version,
                "services_healthy": services_healthy,
                "supported_formats": [".pptx", ".ppt"],
                "status": self.status.value
            },
            metadata={
                "timestamp": datetime.now().isoformat()
            },
            execution_time=0.0,
            memory_used=0
        )
    
    def cleanup(self) -> bool:
        """Clean up any temporary files"""
        try:
            # Clean up temp files if any
            for temp_file in self._temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except:
                    pass
            
            self._temp_files = []
            self.status = ToolStatus.READY
            return True
            
        except Exception as e:
            logger.error(f"Cleanup failed: {e}")
            return False